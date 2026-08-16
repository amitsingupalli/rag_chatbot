"""Document and image ingestion pipeline."""

from __future__ import annotations

import uuid
from pathlib import Path

from llama_index.core import Document, Settings, VectorStoreIndex
from llama_index.core.node_parser import SentenceSplitter
from llama_index.core.storage.storage_context import StorageContext
from llama_index.vector_stores.chroma import ChromaVectorStore

import chromadb

from backend.config import settings
from backend.rag.ocr import SUPPORTED_IMAGE_TYPES, process_image_bytes


class IngestionPipeline:
    def __init__(self) -> None:
        settings.documents_path.mkdir(parents=True, exist_ok=True)
        settings.uploads_path.mkdir(parents=True, exist_ok=True)
        settings.chroma_path.mkdir(parents=True, exist_ok=True)

        from chromadb.api.client import SharedSystemClient
        SharedSystemClient.clear_system_cache()
        self._chroma_client = chromadb.PersistentClient(
            path=str(settings.chroma_path),
            tenant="default_tenant",
            database="default_database",
        )
        self._collection = self._chroma_client.get_or_create_collection("rag_documents")
        self._vector_store = ChromaVectorStore(chroma_collection=self._collection)
        self._storage_context = StorageContext.from_defaults(vector_store=self._vector_store)
        if settings.embedding_provider.lower() == "gemini":
            try:
                from llama_index.embeddings.gemini import GeminiEmbedding
                keys = settings.get_gemini_api_keys
                curr_key = keys[0] if keys else None
                emb_model = settings.gemini_embedding_model
                if not emb_model.startswith("models/"):
                    emb_model = f"models/{emb_model}"
                Settings.embed_model = GeminiEmbedding(
                    model_name=emb_model,
                    api_key=curr_key,
                )
            except Exception:
                from llama_index.embeddings.huggingface import HuggingFaceEmbedding
                Settings.embed_model = HuggingFaceEmbedding(model_name=settings.embedding_model)
        else:
            from llama_index.embeddings.huggingface import HuggingFaceEmbedding
            Settings.embed_model = HuggingFaceEmbedding(model_name=settings.embedding_model)
        self._splitter = SentenceSplitter(
            chunk_size=settings.chunk_size,
            chunk_overlap=settings.chunk_overlap,
        )

    @property
    def index(self) -> VectorStoreIndex:
        return VectorStoreIndex.from_vector_store(
            self._vector_store,
            storage_context=self._storage_context,
        )

    def _build_documents_from_text(
        self, text: str, metadata: dict, doc_id: str | None = None
    ) -> list[Document]:
        doc = Document(text=text, metadata=metadata, id_=doc_id or str(uuid.uuid4()))
        nodes = self._splitter.get_nodes_from_documents([doc])
        return [Document(text=n.get_content(), metadata=n.metadata) for n in nodes]

    def ingest_text_file(
        self, file_path: Path, user_id: str | None = None, conversation_id: str | None = None
    ) -> int:
        text = file_path.read_text(encoding="utf-8", errors="ignore")
        metadata = {
            "source": file_path.name,
            "type": "text",
            "user_id": user_id or "global",
            "conversation_id": conversation_id or "global",
        }
        docs = self._build_documents_from_text(text, metadata)
        VectorStoreIndex.from_documents(
            docs,
            storage_context=self._storage_context,
            show_progress=False,
        )
        return len(docs)

    def ingest_pdf(
        self, file_path: Path, user_id: str | None = None, conversation_id: str | None = None
    ) -> int:
        pages = []
        # 1. Try DoclingReader for highest accuracy layout & table parsing
        try:
            from llama_index.readers.docling import DoclingReader
            reader = DoclingReader()
            docling_docs = reader.load_data(str(file_path))
            for i, d in enumerate(docling_docs):
                if d.text and d.text.strip():
                    pages.append(
                        Document(
                            text=d.text,
                            metadata={
                                "source": file_path.name,
                                "page": i + 1,
                                "type": "docling_pdf",
                                "user_id": user_id or "global",
                                "conversation_id": conversation_id or "global",
                            },
                        )
                    )
        except Exception as docling_err:
            print(f"[Warning] DoclingReader PDF parse failed: {docling_err}. Trying LlamaParse/PyPDF...")

        if not pages and settings.llama_cloud_api_key:
            try:
                from llama_parse import LlamaParse

                parser = LlamaParse(
                    api_key=settings.llama_cloud_api_key,
                    result_type="markdown",
                    verbose=False,
                )
                parsed_docs = parser.load_data(str(file_path))
                for i, pdoc in enumerate(parsed_docs):
                    if pdoc.text and pdoc.text.strip():
                        pages.append(
                            Document(
                                text=pdoc.text,
                                metadata={
                                    "source": file_path.name,
                                    "page": i + 1,
                                    "type": "pdf_llamaparse",
                                    "user_id": user_id or "global",
                                    "conversation_id": conversation_id or "global",
                                },
                            )
                        )
            except Exception as exc:
                print(f"[Warning] LlamaParse failed: {exc}. Falling back to PyPDF.")

        if not pages:
            from pypdf import PdfReader

            reader = PdfReader(str(file_path))
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text() or ""
                if not page_text.strip() and hasattr(page, "images"):
                    ocr_parts = []
                    for img in page.images:
                        try:
                            res = process_image_bytes(img.data)
                            if res.get("ocr_text") and not res["ocr_text"].startswith("[OCR"):
                                ocr_parts.append(res["ocr_text"])
                        except Exception:
                            pass
                    if ocr_parts:
                        page_text = "\n".join(ocr_parts)

                if page_text.strip():
                    pages.append(
                        Document(
                            text=page_text,
                            metadata={
                                "source": file_path.name,
                                "page": i + 1,
                                "type": "pdf",
                                "user_id": user_id or "global",
                                "conversation_id": conversation_id or "global",
                            },
                        )
                    )
        if not pages:
            return 0
        nodes = self._splitter.get_nodes_from_documents(pages)
        docs = [Document(text=n.get_content(), metadata=n.metadata) for n in nodes]
        VectorStoreIndex.from_documents(
            docs,
            storage_context=self._storage_context,
            show_progress=False,
        )
        return len(docs)

    def ingest_image(
        self, file_path: Path, user_id: str | None = None, conversation_id: str | None = None
    ) -> int:
        data = file_path.read_bytes()
        result = process_image_bytes(data)
        text = (
            f"Image: {file_path.name}\n"
            f"Dimensions: {result['width']}x{result['height']}\n"
            f"OCR Text:\n{result['ocr_text']}"
        )
        metadata = {
            "source": file_path.name,
            "type": "image",
            "user_id": user_id or "global",
            "conversation_id": conversation_id or "global",
            "ocr": result["ocr_text"][:500],
        }
        docs = self._build_documents_from_text(text, metadata)
        VectorStoreIndex.from_documents(
            docs,
            storage_context=self._storage_context,
            show_progress=False,
        )
        return len(docs)

    def ingest_pptx(
        self, file_path: Path, user_id: str | None = None, conversation_id: str | None = None
    ) -> int:
        slides = []
        if settings.llama_cloud_api_key:
            try:
                from llama_parse import LlamaParse

                parser = LlamaParse(
                    api_key=settings.llama_cloud_api_key,
                    result_type="markdown",
                    verbose=False,
                )
                parsed_docs = parser.load_data(str(file_path))
                for i, pdoc in enumerate(parsed_docs):
                    if pdoc.text and pdoc.text.strip():
                        slides.append(
                            Document(
                                text=pdoc.text,
                                metadata={
                                    "source": file_path.name,
                                    "page": i + 1,
                                    "type": "pptx_llamaparse",
                                    "user_id": user_id or "global",
                                    "conversation_id": conversation_id or "global",
                                },
                            )
                        )
            except Exception as exc:
                print(f"[Warning] LlamaParse PPTX failed: {exc}. Falling back to python-pptx.")

        if not slides:
            try:
                from pptx import Presentation

                prs = Presentation(str(file_path))
                for i, slide in enumerate(prs.slides):
                    slide_text_parts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slide_text_parts.append(shape.text)
                    slide_text = "\n".join(slide_text_parts)
                    if slide_text.strip():
                        slides.append(
                            Document(
                                text=slide_text,
                                metadata={
                                    "source": file_path.name,
                                    "page": i + 1,
                                    "type": "pptx",
                                    "user_id": user_id or "global",
                                    "conversation_id": conversation_id or "global",
                                },
                            )
                        )
            except Exception as exc:
                print(f"[Warning] python-pptx parsing failed: {exc}")

        if not slides:
            return 0
        nodes = self._splitter.get_nodes_from_documents(slides)
        docs = [Document(text=n.get_content(), metadata=n.metadata) for n in nodes]
        VectorStoreIndex.from_documents(
            docs,
            storage_context=self._storage_context,
            show_progress=False,
        )
        return len(docs)

    def ingest_csv(
        self, file_path: Path, user_id: str | None = None, conversation_id: str | None = None
    ) -> int:
        text = file_path.read_text(encoding="utf-8", errors="ignore")
        metadata = {
            "source": file_path.name,
            "type": "csv",
            "user_id": user_id or "global",
            "conversation_id": conversation_id or "global",
        }
        docs = self._build_documents_from_text(text, metadata)
        VectorStoreIndex.from_documents(
            docs,
            storage_context=self._storage_context,
            show_progress=False,
        )
        return len(docs)

    def ingest_docx(
        self, file_path: Path, user_id: str | None = None, conversation_id: str | None = None
    ) -> int:
        paragraphs = []
        try:
            import docx

            doc = docx.Document(str(file_path))
            for p in doc.paragraphs:
                if p.text and p.text.strip():
                    paragraphs.append(p.text.strip())

            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        paragraphs.append(row_text)

            full_text = "\n\n".join(paragraphs)
            if not full_text.strip():
                return 0

            metadata = {
                "source": file_path.name,
                "type": "docx",
                "user_id": user_id or "global",
                "conversation_id": conversation_id or "global",
            }
            docs = self._build_documents_from_text(full_text, metadata)
            VectorStoreIndex.from_documents(
                docs,
                storage_context=self._storage_context,
                show_progress=False,
            )
            return len(docs)
        except Exception as exc:
            print(f"[Error] docx parsing failed: {exc}")
            return 0

    def ingest_xlsx(
        self, file_path: Path, user_id: str | None = None, conversation_id: str | None = None
    ) -> int:
        sheet_texts = []
        try:
            import pandas as pd

            excel = pd.ExcelFile(str(file_path))
            for sheet_name in excel.sheet_names:
                df = pd.read_excel(excel, sheet_name=sheet_name)
                if not df.empty:
                    csv_str = df.to_csv(index=False)
                    sheet_texts.append(f"--- Sheet: {sheet_name} ---\n{csv_str}")

            full_text = "\n\n".join(sheet_texts)
            if not full_text.strip():
                return 0

            metadata = {
                "source": file_path.name,
                "type": "xlsx",
                "user_id": user_id or "global",
                "conversation_id": conversation_id or "global",
            }
            docs = self._build_documents_from_text(full_text, metadata)
            VectorStoreIndex.from_documents(
                docs,
                storage_context=self._storage_context,
                show_progress=False,
            )
            return len(docs)
        except Exception as exc:
            print(f"[Error] xlsx parsing failed: {exc}")
            return 0

    def ingest_html(
        self, file_path: Path, user_id: str | None = None, conversation_id: str | None = None
    ) -> int:
        try:
            from bs4 import BeautifulSoup

            raw_html = file_path.read_text(encoding="utf-8", errors="ignore")
            soup = BeautifulSoup(raw_html, "html.parser")
            for script in soup(["script", "style"]):
                script.decompose()

            text = soup.get_text(separator="\n")
            lines = (line.strip() for line in text.splitlines())
            clean_text = "\n".join(chunk for chunk in lines if chunk)

            if not clean_text.strip():
                return 0

            metadata = {
                "source": file_path.name,
                "type": "html",
                "user_id": user_id or "global",
                "conversation_id": conversation_id or "global",
            }
            docs = self._build_documents_from_text(clean_text, metadata)
            VectorStoreIndex.from_documents(
                docs,
                storage_context=self._storage_context,
                show_progress=False,
            )
            return len(docs)
        except Exception as exc:
            print(f"[Error] html parsing failed: {exc}")
            return 0

    def ingest_file(
        self, file_path: str | Path, user_id: str | None = None, conversation_id: str | None = None
    ) -> int:
        file_path = Path(file_path)
        suffix = file_path.suffix.lower()
        if suffix in SUPPORTED_IMAGE_TYPES:
            return self.ingest_image(file_path, user_id, conversation_id)
        if suffix == ".pdf":
            return self.ingest_pdf(file_path, user_id, conversation_id)
        if suffix in {".doc", ".docx"}:
            return self.ingest_docx(file_path, user_id, conversation_id)
        if suffix in {".ppt", ".pptx"}:
            return self.ingest_pptx(file_path, user_id, conversation_id)
        if suffix in {".xls", ".xlsx"}:
            return self.ingest_xlsx(file_path, user_id, conversation_id)
        if suffix in {".html", ".htm"}:
            return self.ingest_html(file_path, user_id, conversation_id)
        if suffix in {".txt", ".md", ".csv", ".json", ".tsv"}:
            return (
                self.ingest_csv(file_path, user_id, conversation_id)
                if suffix == ".csv"
                else self.ingest_text_file(file_path, user_id, conversation_id)
            )
        raise ValueError(f"Unsupported file type: {suffix}")

    def ingest_bytes(
        self,
        data: bytes,
        filename: str,
        user_id: str | None = None,
        conversation_id: str | None = None,
    ) -> int:
        import re
        p = Path(filename)
        safe_stem = re.sub(r"[^\w\-]", "_", p.stem)
        safe_ext = p.suffix
        safe_filename = f"{safe_stem}{safe_ext}"
        dest = settings.uploads_path / f"{uuid.uuid4().hex}_{safe_filename}"
        dest.write_bytes(data)
        return self.ingest_file(dest, user_id, conversation_id)
